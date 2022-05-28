import os

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

# Change the current working directory.
os.chdir("C:\\Users\\jegat\\Desktop\\ppt\\pictures")
# Store images directory.
pictures = os.listdir()

# Create a presentation object.
ppt = Presentation()
# Selecting blank slides.
blank_slide = ppt.slide_layouts[6]

height = Inches(7.5)
width = Inches(11)
left = Inches(-0.8)
top = Inches(0)

# Process the images.
for i in range(len(pictures)):
    picture = Image.open(pictures[i])
    cropped = picture.crop((148, 0, 1414, 720))
    cropped.save(f"{i}.jpg")
    # Delete old pictures.
    os.remove(pictures[i])

    # Attaching slides to ppt.
    slide = ppt.slides.add_slide(blank_slide)

    # Adding images.
    pic = slide.shapes.add_picture(f"{i}.jpg", left, top, height=height, width=width)

ppt.save("bot.pptx")
